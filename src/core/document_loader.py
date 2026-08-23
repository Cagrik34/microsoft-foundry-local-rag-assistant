"""
Belge İşleme ve Metin Çıkarma Modülü (src/core/document_loader.py)
===================================================================
PDF, DOCX, XLSX, PPTX, TXT ve MD dosyalarından yapı duyarlı (structure-aware) metin ayıklar ve öbekler.
Excel tablolarını Markdown tablo formatına dönüştürerek satır/sütun ilişkisini korur.
"""

import os
import re
from typing import List
from src.config import DOCUMENTS_DIR, SUPPORTED_EXTENSIONS, CHUNK_SIZE, CHUNK_OVERLAP, MIN_CHUNK_LENGTH
from src.core.models import TextChunk, DocumentInfo

# Opsiyonel belge okuyucu kütüphaneleri
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pptx
except ImportError:
    pptx = None


def scan_documents(directory: str = DOCUMENTS_DIR) -> List[str]:
    """Dizindeki desteklenen uzantıya sahip tüm belgeleri tarar ve yollarını döndürür."""
    if not os.path.isdir(directory):
        print(f"⚠️  Belge dizini bulunamadı: {directory}")
        return []

    found_files: List[str] = []
    for root, _dirs, files in os.walk(directory):
        for file_name in sorted(files):
            ext = os.path.splitext(file_name)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                found_files.append(os.path.join(root, file_name))
    return found_files


def _extract_pdf(file_path: str) -> str:
    """PDF dosyasından sayfa sayfa metin çıkarır."""
    if pypdf is None:
        print("⚠️  PDF okunamadı: pypdf kütüphanesi yüklü değil.")
        return ""
    try:
        reader = pypdf.PdfReader(file_path)
        parts = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                parts.append(f"=== Sayfa {i+1} ===\n{text.strip()}")
        return "\n\n".join(parts)
    except Exception as e:
        print(f"⚠️  PDF okuma hatası ({file_path}): {e}")
        return ""


def _extract_docx(file_path: str) -> str:
    """Word (.docx) dosyasından başlık hiyerarşisi ve tabloları Markdown formatında çıkarır."""
    if docx is None:
        print("⚠️  DOCX okunamadı: python-docx kütüphanesi yüklü değil.")
        return ""
    try:
        doc = docx.Document(file_path)
        parts = []
        for elem in doc.paragraphs:
            txt = elem.text.strip()
            if not txt:
                continue
            if elem.style.name.startswith("Heading 1"):
                parts.append(f"\n# {txt}\n")
            elif elem.style.name.startswith("Heading 2"):
                parts.append(f"\n## {txt}\n")
            elif elem.style.name.startswith("Heading 3"):
                parts.append(f"\n### {txt}\n")
            else:
                parts.append(txt)

        # Tabloları Markdown formatına dönüştür
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                table_rows.append(cells)

            if table_rows:
                header = table_rows[0]
                md_table = [
                    "| " + " | ".join(header) + " |",
                    "| " + " | ".join(["---"] * len(header)) + " |"
                ]
                for r in table_rows[1:]:
                    # Sütun sayısını eşitle
                    padded = r + [""] * (len(header) - len(r))
                    md_table.append("| " + " | ".join(padded[:len(header)]) + " |")
                parts.append("\n" + "\n".join(md_table) + "\n")

        return "\n\n".join(parts)
    except Exception as e:
        print(f"⚠️  DOCX okuma hatası ({file_path}): {e}")
        return ""


def _extract_xlsx(file_path: str) -> str:
    """Excel (.xlsx) sayfalarını Markdown Tablo formatına dönüştürerek sütun/satır ilişkisini korur."""
    if openpyxl is None:
        print("⚠️  XLSX okunamadı: openpyxl kütüphanesi yüklü değil.")
        return ""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        sheet_parts = []
        for name in wb.sheetnames:
            rows = list(wb[name].iter_rows(values_only=True))
            if not rows:
                continue

            # Boş satırları filtrele
            valid_rows = []
            for r in rows:
                cleaned = [str(v).strip() if v is not None else "" for v in r]
                if any(cleaned):
                    valid_rows.append(cleaned)

            if not valid_rows:
                continue

            header = valid_rows[0]
            col_count = len(header)
            md_lines = [
                f"\n=== Sayfa: {name} ===\n",
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(["---"] * col_count) + " |"
            ]
            for r in valid_rows[1:]:
                padded = r + [""] * (col_count - len(r))
                md_lines.append("| " + " | ".join(padded[:col_count]) + " |")

            sheet_parts.append("\n".join(md_lines))

        wb.close()
        return "\n\n".join(sheet_parts)
    except Exception as e:
        print(f"⚠️  XLSX okuma hatası ({file_path}): {e}")
        return ""


def _extract_pptx(file_path: str) -> str:
    """PowerPoint (.pptx) sunumundaki slayt metinlerini çıkarır."""
    if pptx is None:
        print("⚠️  PPTX okunamadı: python-pptx kütüphanesi yüklü değil.")
        return ""
    try:
        prs = pptx.Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"=== Slayt {i+1} ===")
            slide_txt = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if slide_txt:
                parts.append("\n".join(slide_txt))
        return "\n\n".join(parts)
    except Exception as e:
        print(f"⚠️  PPTX okuma hatası ({file_path}): {e}")
        return ""


def read_document(file_path: str) -> DocumentInfo:
    """Dosya uzantısına göre uygun okuyucuyu çağırır ve metni döndürür."""
    file_name = os.path.basename(file_path)
    ext = os.path.splitext(file_name)[1].lower()

    if ext in (".md", ".txt"):
        content = ""
        for enc in ["utf-8", "utf-8-sig", "cp1254", "latin-1"]:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    content = f.read()
                break
            except (UnicodeDecodeError, UnicodeError):
                continue
            except Exception as e:
                print(f"⚠️  Metin dosyası okuma hatası ({file_path}): {e}")
                content = ""
                break
    elif ext == ".pdf":
        content = _extract_pdf(file_path)
    elif ext == ".docx":
        content = _extract_docx(file_path)
    elif ext == ".xlsx":
        content = _extract_xlsx(file_path)
    elif ext == ".pptx":
        content = _extract_pptx(file_path)
    else:
        content = ""

    return DocumentInfo(file_path=file_path, file_name=file_name, content=content)


def chunk_text(text: str) -> List[str]:
    """Yapı ve başlık duyarlı semantik öbekleme algoritması."""
    if not text or not text.strip():
        return []

    # Bölüm ve başlık sınırlarına göre böl
    sections = [s.strip() for s in re.split(r'\n(?=== |#+ )', text) if s.strip()]
    chunks: List[str] = []
    current = ""

    for sec in sections:
        if len(sec) > CHUNK_SIZE:
            if current:
                chunks.append(current)
                current = ""

            paragraphs = [p.strip() for p in sec.split("\n\n") if p.strip()]
            p_current = ""
            for p in paragraphs:
                if len(p) > CHUNK_SIZE:
                    if p_current:
                        chunks.append(p_current)
                        p_current = ""
                    # Cümle bazlı kayan pencere
                    start = 0
                    while start < len(p):
                        end = start + CHUNK_SIZE
                        split = p.rfind(". ", start, end) if end < len(p) else len(p)
                        if split <= start:
                            split = p.rfind(" ", start, end) if end < len(p) else len(p)
                        if split <= start:
                            split = end
                        else:
                            split += 1  # Noktayı dahil et
                        chunks.append(p[start:split].strip())
                        start = max(0, split - CHUNK_OVERLAP) if split < len(p) else len(p)
                    continue

                test_p = f"{p_current}\n\n{p}" if p_current else p
                if len(test_p) <= CHUNK_SIZE:
                    p_current = test_p
                else:
                    if p_current:
                        chunks.append(p_current)
                    p_current = p

            if p_current:
                chunks.append(p_current)
            continue

        test_sec = f"{current}\n\n{sec}" if current else sec
        if len(test_sec) <= CHUNK_SIZE:
            current = test_sec
        else:
            if current:
                chunks.append(current)
            current = sec

    if current:
        chunks.append(current)

    return [c.strip() for c in chunks if len(c.strip()) >= MIN_CHUNK_LENGTH]


def process_document(file_path: str) -> DocumentInfo:
    """Tek bir belgeyi okur ve metin öbeklerini (chunks) oluşturur."""
    doc = read_document(file_path)
    if not doc.content or not doc.content.strip():
        doc.chunks = []
        return doc
    raw_chunks = chunk_text(doc.content)
    doc.chunks = [
        TextChunk(content=c, source_file=doc.file_name, chunk_index=i)
        for i, c in enumerate(raw_chunks)
    ]
    return doc


def process_all_documents(directory: str = DOCUMENTS_DIR) -> List[DocumentInfo]:
    """Taranan tüm belgeleri okur ve öbeklerine ayırır."""
    paths = scan_documents(directory)
    docs = []
    for p in paths:
        try:
            d = process_document(p)
            docs.append(d)
            print(f"  ✅ {d.file_name} → {len(d.chunks)} öbek")
        except Exception as e:
            print(f"  ❌ {os.path.basename(p)} → Hata: {e}")
    return docs
